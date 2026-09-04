
import io
import os
import re
import hashlib
from typing import List, Dict, Any

import numpy as np
import pandas as pd
import streamlit as st
from openai import OpenAI
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


# ============================================================
# PAGE CONFIG
# ============================================================
st.set_page_config(
    page_title="DocuMind AI",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ============================================================
# CUSTOM CSS — keeps the visual language of the supplied HTML
# ============================================================
st.markdown(
    """
<style>
@import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&display=swap');

:root {
    --indigo: #635BFF;
    --indigo-dark: #534be7;
    --text: #0f172a;
    --muted: #64748b;
    --border: #e2e8f0;
    --soft: #f8fafc;
}

html, body, [class*="css"] {
    font-family: 'Plus Jakarta Sans', sans-serif;
}

.stApp {
    background: #fbfbfe;
    color: var(--text);
}

#MainMenu, footer, header { visibility: hidden; }

.block-container {
    max-width: 1440px;
    padding: 1.25rem 1.4rem 2rem;
}

section[data-testid="stSidebar"] {
    background: #ffffff;
    border-right: 1px solid rgba(226,232,240,.85);
}

section[data-testid="stSidebar"] > div {
    padding: 1.25rem;
}

div[data-testid="stFileUploader"] {
    background: #fff;
    border: 1px dashed #cbd5e1;
    border-radius: 16px;
    padding: 10px;
}

div[data-testid="stFileUploader"] section {
    border: 0;
}

.stButton > button {
    border-radius: 11px;
    border: 1px solid #e2e8f0;
    font-family: 'Plus Jakarta Sans', sans-serif;
    font-weight: 600;
    transition: .18s ease;
}

.stButton > button:hover {
    border-color: #c7d2fe;
    transform: translateY(-1px);
}

div[data-testid="stTextInput"] input,
div[data-testid="stTextArea"] textarea,
div[data-testid="stSelectbox"] div[data-baseweb="select"] > div {
    border-radius: 11px;
    border-color: #dbe2ea;
    background: #f8fafc;
}

div[data-testid="stTextInput"] input:focus {
    border-color: #818cf8;
    box-shadow: 0 0 0 2px rgba(99,91,255,.10);
}

.dm-auth-page {
    min-height: 82vh;
    padding: 28px 42px 38px;
    border: 1px solid rgba(255,255,255,.85);
    border-radius: 30px;
    background:
        radial-gradient(circle at 15% 15%, #f1f4ff 0%, #f8faff 45%, #eff2fe 100%);
    box-shadow:
        0 25px 60px -15px rgba(67,56,202,.08),
        0 0 0 1px rgba(255,255,255,.8) inset;
}

.dm-topbar {
    display:flex;
    justify-content:space-between;
    align-items:center;
    gap:20px;
    margin-bottom:45px;
}

.dm-brand {
    display:flex;
    align-items:center;
    gap:13px;
}

.dm-logo {
    width:48px;
    height:48px;
    border-radius:16px;
    background:linear-gradient(135deg,#2563eb,#6366f1,#8b5cf6);
    display:flex;
    align-items:center;
    justify-content:center;
    color:#fff;
    font-size:22px;
    font-weight:800;
    box-shadow:0 10px 25px rgba(99,91,255,.20);
}

.dm-brand-name {
    font-size:20px;
    line-height:1;
    font-weight:800;
    letter-spacing:-.6px;
}

.dm-brand-name span { color:var(--indigo); }

.dm-brand-sub {
    color:#94a3b8;
    font-size:11px;
    margin-top:5px;
    font-weight:500;
}

.dm-trust {
    color:#64748b;
    font-size:11px;
    font-weight:500;
}

.dm-hero-title {
    font-size:clamp(40px,4.3vw,58px);
    line-height:1.12;
    letter-spacing:-2.6px;
    font-weight:800;
    margin:0 0 18px;
}

.dm-gradient {
    background:linear-gradient(135deg,#2563eb 0%,#6366f1 70%,#8b5cf6 100%);
    -webkit-background-clip:text;
    background-clip:text;
    -webkit-text-fill-color:transparent;
}

.dm-hero-copy {
    color:#64748b;
    font-size:15px;
    line-height:1.75;
    max-width:620px;
    margin-bottom:27px;
}

.dm-badge {
    display:inline-block;
    padding:7px 13px;
    border-radius:999px;
    color:#4f46e5;
    background:#eef2ff;
    border:1px solid #e0e7ff;
    font-size:10px;
    font-weight:800;
    letter-spacing:1px;
    text-transform:uppercase;
    margin-bottom:20px;
}

.dm-features {
    display:grid;
    grid-template-columns:repeat(4,1fr);
    gap:10px;
    max-width:640px;
    margin-bottom:28px;
}

.dm-feature {
    display:flex;
    flex-direction:column;
    gap:8px;
    color:#334155;
    font-size:11px;
    font-weight:700;
}

.dm-feature-icon {
    width:38px;
    height:38px;
    border-radius:12px;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:16px;
}

.dm-showcase {
    position:relative;
    height:245px;
    max-width:610px;
    margin-top:8px;
}

.dm-arc {
    position:absolute;
    left:5%;
    right:5%;
    top:8px;
    height:150px;
    border:1px dashed #c7d2fe;
    border-radius:50%;
}

.dm-doc {
    position:absolute;
    width:105px;
    height:145px;
    border-radius:12px;
    background:#fff;
    border:1px solid #e2e8f0;
    box-shadow:0 12px 28px -6px rgba(45,55,72,.12);
    display:flex;
    flex-direction:column;
    align-items:center;
    justify-content:center;
    gap:10px;
}

.dm-doc .mini-icon {
    width:42px;
    height:50px;
    border-radius:8px;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:24px;
}

.dm-doc-center {
    left:50%;
    top:30px;
    transform:translateX(-50%);
    width:110px;
    height:155px;
    z-index:4;
    box-shadow:0 20px 35px -8px rgba(79,70,229,.22);
}

.dm-doc-left { left:22%; top:62px; transform:rotate(-10deg); z-index:2; }
.dm-doc-right { right:22%; top:62px; transform:rotate(10deg); z-index:2; }

.dm-doc-far-left { left:8%; top:82px; transform:rotate(-17deg); z-index:1; }
.dm-doc-far-right { right:8%; top:82px; transform:rotate(17deg); z-index:1; }

.dm-tagline {
    position:absolute;
    left:12%;
    right:12%;
    bottom:5px;
    text-align:center;
    color:#64748b;
    font-size:12px;
    font-style:italic;
    font-weight:600;
}

.dm-auth-card {
    background:#fff;
    border:1px solid #eef2f7;
    border-radius:30px;
    padding:35px;
    box-shadow:0 20px 40px -15px rgba(0,15,60,.08);
    max-width:430px;
    margin:0 auto;
}

.dm-lock {
    width:56px;
    height:56px;
    border-radius:17px;
    display:flex;
    align-items:center;
    justify-content:center;
    background:#eef2ff;
    color:#4f46e5;
    margin:0 auto 18px;
    font-size:22px;
}

.dm-auth-title {
    text-align:center;
    font-size:25px;
    font-weight:800;
    letter-spacing:-.8px;
    margin-bottom:8px;
}

.dm-auth-copy {
    color:#64748b;
    font-size:12px;
    line-height:1.7;
    text-align:center;
    margin:0 auto 22px;
    max-width:310px;
}

.dm-secure {
    text-align:center;
    color:#94a3b8;
    font-size:10px;
    margin-top:8px;
}

.dm-workspace-head {
    display:flex;
    align-items:center;
    justify-content:space-between;
    gap:18px;
    padding:3px 2px 18px;
}

.dm-workspace-brand {
    display:flex;
    align-items:center;
    gap:10px;
}

.dm-small-logo {
    width:32px;
    height:32px;
    border-radius:9px;
    background:#635BFF;
    color:#fff;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:15px;
}

.dm-workspace-name {
    font-size:17px;
    font-weight:800;
    letter-spacing:-.5px;
}

.dm-workspace-sub {
    color:#94a3b8;
    font-size:10px;
    margin-top:2px;
}

.dm-workspace-trust {
    color:#94a3b8;
    font-size:10px;
}

.dm-pill {
    display:inline-block;
    padding:6px 11px;
    border-radius:999px;
    color:#4f46e5;
    background:#f3f1ff;
    border:1px solid #e7e3ff;
    font-size:10px;
    font-weight:700;
    text-transform:uppercase;
    letter-spacing:.7px;
}

.dm-center-title {
    text-align:center;
    font-size:35px;
    font-weight:800;
    letter-spacing:-1.6px;
    margin:12px 0 5px;
}

.dm-center-copy {
    text-align:center;
    color:#64748b;
    font-size:12px;
    margin-bottom:18px;
}

.dm-banner {
    background:#ebf9f1;
    border:1px solid #bff0d4;
    border-radius:11px;
    padding:9px 13px;
    color:#166534;
    font-size:10px;
    font-weight:600;
    margin-bottom:15px;
}

.dm-stat-grid {
    display:grid;
    grid-template-columns:repeat(4,1fr);
    gap:10px;
    margin:12px 0 16px;
}

.dm-stat {
    background:#fff;
    border:1px solid #e2e8f0;
    border-radius:15px;
    padding:13px;
    display:flex;
    align-items:center;
    gap:10px;
    box-shadow:0 3px 10px rgba(15,23,42,.025);
}

.dm-stat-icon {
    width:39px;
    height:39px;
    border-radius:11px;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:15px;
    flex:0 0 auto;
}

.dm-stat-number { font-size:19px; font-weight:800; line-height:1; }
.dm-stat-label { font-size:9px; color:#94a3b8; margin-top:4px; }
.dm-stat-status { font-size:13px; font-weight:800; line-height:1; }

.dm-file-row {
    background:#fff;
    border:1px solid #e2e8f0;
    border-radius:11px;
    padding:10px 12px;
    display:flex;
    align-items:center;
    justify-content:space-between;
    gap:12px;
    margin-bottom:8px;
}

.dm-file-main {
    display:flex;
    align-items:center;
    gap:8px;
    min-width:0;
}

.dm-file-name {
    font-size:10px;
    font-weight:700;
    color:#334155;
    white-space:nowrap;
    overflow:hidden;
    text-overflow:ellipsis;
}

.dm-file-meta {
    color:#94a3b8;
    font-size:9px;
}

.dm-file-status {
    color:#16a34a;
    font-size:9px;
    font-weight:700;
    white-space:nowrap;
}

.dm-question {
    background:#fff;
    border:1px solid #e2e8f0;
    border-radius:11px;
    padding:12px;
    display:flex;
    align-items:center;
    gap:10px;
    margin-top:9px;
}

.dm-avatar-user {
    width:27px;
    height:27px;
    border-radius:50%;
    background:#ef4444;
    color:#fff;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:12px;
    flex:0 0 auto;
}

.dm-answer {
    background:#fff;
    border:1px solid #e2e8f0;
    border-radius:11px;
    padding:13px;
    margin-top:8px;
}

.dm-avatar-ai {
    width:27px;
    height:27px;
    border-radius:50%;
    background:#f59e0b;
    color:#fff;
    display:flex;
    align-items:center;
    justify-content:center;
    font-size:12px;
    flex:0 0 auto;
}

.dm-answer-head {
    display:flex;
    align-items:center;
    gap:9px;
    font-size:10px;
    font-weight:700;
    margin-bottom:8px;
}

.dm-answer-body {
    color:#64748b;
    font-size:10px;
    line-height:1.65;
}

.dm-sources {
    border-top:1px solid #f1f5f9;
    margin-top:10px;
    padding-top:9px;
    color:#64748b;
    font-size:9px;
}

.dm-input-wrap {
    background:#fff;
    border:1px solid #c7d2fe;
    border-radius:15px;
    padding:5px 7px 5px 10px;
    box-shadow:0 4px 15px rgba(99,91,255,.06);
}

.dm-disclaimer {
    text-align:center;
    color:#94a3b8;
    font-size:8px;
    margin-top:6px;
}

.dm-empty {
    border:1px dashed #cbd5e1;
    background:rgba(255,255,255,.7);
    border-radius:15px;
    padding:25px;
    text-align:center;
    color:#64748b;
}

.dm-empty-icon {
    font-size:28px;
    margin-bottom:7px;
}

@media (max-width: 900px) {
    .dm-auth-page { padding:22px; }
    .dm-topbar { margin-bottom:25px; }
    .dm-trust { display:none; }
    .dm-features { grid-template-columns:repeat(2,1fr); }
    .dm-showcase { transform:scale(.88); transform-origin:top center; margin-bottom:-20px; }
    .dm-stat-grid { grid-template-columns:repeat(2,1fr); }
}

@media (max-width: 620px) {
    .block-container { padding: .6rem .7rem 1.2rem; }
    .dm-auth-page { padding:18px; border-radius:22px; }
    .dm-hero-title { font-size:38px; letter-spacing:-1.8px; }
    .dm-features { grid-template-columns:repeat(2,1fr); }
    .dm-showcase { display:none; }
    .dm-auth-card { padding:25px 20px; border-radius:23px; }
    .dm-center-title { font-size:28px; }
    .dm-workspace-trust { display:none; }
}
</style>
""",
    unsafe_allow_html=True,
)


# ============================================================
# SESSION STATE
# ============================================================
defaults = {
    "connected": False,
    "api_key": "",
    "model": "gpt-4o-mini",
    "client": None,
    "documents": [],
    "chunks": [],
    "embeddings": None,
    "messages": [],
    "settings": {
        "chunk_size": 500,
        "chunk_overlap": 50,
        "retrieved_chunks": 3,
        "temperature": 0.2,
        "embedding_model": "text-embedding-3-small",
    },
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v


# ============================================================
# HELPERS
# ============================================================
def file_icon(name: str) -> str:
    ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""
    return {
        "pdf": "📕",
        "docx": "📘",
        "doc": "📘",
        "xlsx": "📊",
        "xls": "📊",
        "csv": "📊",
        "pptx": "📙",
        "ppt": "📙",
        "txt": "📄",
        "md": "📄",
        "png": "🖼️",
        "jpg": "🖼️",
        "jpeg": "🖼️",
    }.get(ext, "📄")


def human_size(n: int) -> str:
    if n < 1024:
        return f"{n} B"
    if n < 1024**2:
        return f"{n/1024:.1f} KB"
    return f"{n/1024**2:.1f} MB"


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def chunk_text(text: str, size: int, overlap: int) -> List[str]:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return []

    size = max(100, int(size))
    overlap = min(max(0, int(overlap)), size - 1)

    result = []
    start = 0
    while start < len(text):
        end = min(len(text), start + size)
        part = text[start:end].strip()
        if part:
            result.append(part)
        if end >= len(text):
            break
        start = end - overlap
    return result


def parse_pdf(data: bytes) -> List[Dict[str, Any]]:
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"text": text, "page": i})
    return pages


def parse_docx(data: bytes) -> List[Dict[str, Any]]:
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            text += "\n" + "\n".join(rows)
    return [{"text": text, "page": None}] if text.strip() else []


def parse_pptx(data: bytes) -> List[Dict[str, Any]]:
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        text = "\n".join(texts)
        if text.strip():
            slides.append({"text": text, "slide": i})
    return slides


def parse_excel(data: bytes, name: str) -> List[Dict[str, Any]]:
    xls = pd.ExcelFile(io.BytesIO(data))
    result = []
    for sheet in xls.sheet_names:
        df = pd.read_excel(io.BytesIO(data), sheet_name=sheet)
        if df.empty:
            continue
        lines = [f"Sheet: {sheet}"]
        for idx, row in df.iterrows():
            vals = []
            for col, value in row.items():
                if pd.notna(value):
                    vals.append(f"{col}: {value}")
            if vals:
                lines.append(f"Row {idx + 2}: " + " | ".join(vals))
        result.append({"text": "\n".join(lines), "sheet": sheet})
    return result


def parse_csv(data: bytes) -> List[Dict[str, Any]]:
    df = pd.read_csv(io.BytesIO(data))
    lines = []
    for idx, row in df.iterrows():
        vals = [f"{col}: {value}" for col, value in row.items() if pd.notna(value)]
        if vals:
            lines.append(f"Row {idx + 2}: " + " | ".join(vals))
    return [{"text": "\n".join(lines), "page": None}] if lines else []


def parse_text(data: bytes) -> List[Dict[str, Any]]:
    text = data.decode("utf-8", errors="ignore")
    return [{"text": text, "page": None}] if text.strip() else []


def parse_uploaded_file(uploaded) -> List[Dict[str, Any]]:
    name = uploaded.name
    data = uploaded.getvalue()
    ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""

    if ext == "pdf":
        return parse_pdf(data)
    if ext in {"docx", "doc"}:
        return parse_docx(data)
    if ext in {"pptx", "ppt"}:
        return parse_pptx(data)
    if ext in {"xlsx", "xls"}:
        return parse_excel(data, name)
    if ext == "csv":
        return parse_csv(data)
    if ext in {"txt", "md"}:
        return parse_text(data)

    # Image OCR is intentionally optional so Streamlit Cloud remains easy
    # to deploy without a system-level Tesseract dependency.
    return []


def build_chunks(files: List[Any]) -> List[Dict[str, Any]]:
    all_chunks = []
    size = st.session_state.settings["chunk_size"]
    overlap = st.session_state.settings["chunk_overlap"]

    for uploaded in files:
        raw = uploaded.getvalue()
        signature = sha256_bytes(raw)
        parsed = parse_uploaded_file(uploaded)

        for item in parsed:
            pieces = chunk_text(item.get("text", ""), size, overlap)
            for piece in pieces:
                meta = {
                    "filename": uploaded.name,
                    "file_type": uploaded.name.rsplit(".", 1)[-1].upper()
                    if "." in uploaded.name else "FILE",
                    "page": item.get("page"),
                    "sheet": item.get("sheet"),
                    "slide": item.get("slide"),
                    "signature": signature,
                    "preview": piece[:500],
                }
                all_chunks.append({"text": piece, "metadata": meta})
    return all_chunks


def get_embeddings(client: OpenAI, texts: List[str]) -> np.ndarray:
    model = st.session_state.settings["embedding_model"]
    response = client.embeddings.create(model=model, input=texts)
    return np.asarray([item.embedding for item in response.data], dtype=np.float32)


def retrieve(question: str, top_k: int) -> List[Dict[str, Any]]:
    if not st.session_state.chunks or st.session_state.embeddings is None:
        return []

    client = st.session_state.client
    q = get_embeddings(client, [question])[0]

    doc_vectors = st.session_state.embeddings
    q_norm = np.linalg.norm(q) or 1.0
    d_norm = np.linalg.norm(doc_vectors, axis=1)
    scores = (doc_vectors @ q) / (d_norm * q_norm + 1e-8)

    indices = np.argsort(scores)[::-1][:top_k]
    results = []
    for idx in indices:
        item = dict(st.session_state.chunks[idx])
        item["score"] = float(scores[idx])
        results.append(item)
    return results


def answer_question(question: str) -> tuple[str, List[Dict[str, Any]]]:
    sources = retrieve(question, st.session_state.settings["retrieved_chunks"])

    if not sources:
        return (
            "I don't know based on the uploaded documents. "
            "Upload and index a supported document first.",
            [],
        )

    context_parts = []
    for i, source in enumerate(sources, start=1):
        m = source["metadata"]
        location = ""
        if m.get("page"):
            location = f"Page {m['page']}"
        elif m.get("sheet"):
            location = f"Sheet: {m['sheet']}"
        elif m.get("slide"):
            location = f"Slide {m['slide']}"
        context_parts.append(
            f"[Source {i}] {m['filename']} {location}\n{source['text']}"
        )

    context = "\n\n".join(context_parts)

    history = []
    for msg in st.session_state.messages[-6:]:
        history.append({"role": msg["role"], "content": msg["content"]})

    system = """You are DocuMind AI, a document-grounded assistant.
Answer using ONLY the supplied document context.
If the answer is not supported by the context, say you don't know based on the uploaded documents.
Do not invent facts, citations, pages, values, or sources.
Keep the answer concise and useful.
"""

    response = st.session_state.client.chat.completions.create(
        model=st.session_state.model,
        temperature=float(st.session_state.settings["temperature"]),
        messages=[
            {"role": "system", "content": system},
            *history,
            {
                "role": "user",
                "content": f"Document context:\n\n{context}\n\nQuestion: {question}",
            },
        ],
    )
    return response.choices[0].message.content.strip(), sources


def reset_workspace():
    st.session_state.documents = []
    st.session_state.chunks = []
    st.session_state.embeddings = None
    st.session_state.messages = []


# ============================================================
# SCREEN 1 — CONNECTION / LANDING PAGE
# ============================================================
def render_connection():
    st.markdown('<div class="dm-auth-page">', unsafe_allow_html=True)

    st.markdown(
        """
<div class="dm-topbar">
  <div class="dm-brand">
    <div class="dm-logo">▤</div>
    <div>
      <div class="dm-brand-name">DocuMind <span>AI</span></div>
      <div class="dm-brand-sub">Your documents. Now searchable with AI.</div>
    </div>
  </div>
  <div class="dm-trust">♙ Secure &nbsp; • &nbsp; Private &nbsp; • &nbsp; Powered by OpenAI</div>
</div>
""",
        unsafe_allow_html=True,
    )

    left, right = st.columns([1.55, 1], gap="large")

    with left:
        st.markdown(
            """
<div class="dm-badge">AI Document Intelligence</div>
<div class="dm-hero-title">
  Your documents.<br>
  <span class="dm-gradient">Now searchable with AI.</span>
</div>
<div class="dm-hero-copy">
  Upload documents, spreadsheets, presentations and images, then ask
  questions using natural language.
</div>

<div class="dm-features">
  <div class="dm-feature">
    <div class="dm-feature-icon" style="background:#eef2ff;color:#4f46e5">▤</div>
    All File Types
  </div>
  <div class="dm-feature">
    <div class="dm-feature-icon" style="background:#ecfdf5;color:#10b981">ϟ</div>
    Instant Answers
  </div>
  <div class="dm-feature">
    <div class="dm-feature-icon" style="background:#fff7ed;color:#f97316">✓</div>
    100% Private
  </div>
  <div class="dm-feature">
    <div class="dm-feature-icon" style="background:#eff6ff;color:#3b82f6">⚙</div>
    Powered by OpenAI
  </div>
</div>

<div class="dm-showcase">
  <div class="dm-arc"></div>
  <div class="dm-doc dm-doc-far-left"><div class="mini-icon" style="background:#fff1f2">📕</div><small>PDF</small></div>
  <div class="dm-doc dm-doc-left"><div class="mini-icon" style="background:#ecfdf5">📊</div><small>XLSX</small></div>
  <div class="dm-doc dm-doc-center"><div class="mini-icon" style="background:#eef2ff">▤</div><small style="font-weight:700">DOCS</small></div>
  <div class="dm-doc dm-doc-right"><div class="mini-icon" style="background:#eff6ff">📘</div><small>DOCX</small></div>
  <div class="dm-doc dm-doc-far-right"><div class="mini-icon" style="background:#fff7ed">📙</div><small>PPTX</small></div>
  <div class="dm-tagline">“Turn your documents into knowledge.”</div>
</div>
""",
            unsafe_allow_html=True,
        )

    with right:
        st.markdown(
            """
<div class="dm-auth-card">
  <div class="dm-lock">♙</div>
  <div class="dm-auth-title">Connect to OpenAI</div>
  <div class="dm-auth-copy">
    Enter your OpenAI API key and select a model to get started.
    Your API key is kept only in this session and is not stored.
  </div>
</div>
""",
            unsafe_allow_html=True,
        )

        # The inputs are deliberately outside the HTML card so Streamlit can
        # handle them natively and securely.
        api_key = st.text_input(
            "OpenAI API Key",
            type="password",
            placeholder="sk-...",
            label_visibility="visible",
            key="connect_api_key",
        )

        model = st.selectbox(
            "Select Model",
            ["gpt-4o-mini", "gpt-4.1-mini", "gpt-4o"],
            index=0,
            key="connect_model",
        )

        st.caption("✓ Your key is secure and never stored on our servers.")

        if st.button("⚡  Connect to OpenAI  →", use_container_width=True, type="primary"):
            if not api_key.strip():
                st.error("Please enter your OpenAI API key.")
            else:
                with st.spinner("Connecting securely..."):
                    try:
                        client = OpenAI(api_key=api_key.strip())
                        client.models.list()
                        st.session_state.api_key = api_key.strip()
                        st.session_state.model = model
                        st.session_state.client = client
                        st.session_state.connected = True
                        st.rerun()
                    except Exception:
                        st.error(
                            "We couldn't connect with that API key. "
                            "Please check the key and try again."
                        )

        st.markdown(
            """
<div style="text-align:center;margin:16px 0 8px;color:#94a3b8;font-size:10px">
  ──────────────── &nbsp; New to OpenAI? &nbsp; ────────────────
</div>
""",
            unsafe_allow_html=True,
        )

        st.link_button(
            "Get your API key ↗",
            "https://platform.openai.com/api-keys",
            use_container_width=False,
        )

        st.markdown(
            """
<div class="dm-secure">
  🔒 Secure Session Only &nbsp;&nbsp; ◉ Never Stored &nbsp;&nbsp; ✓ You Control Your Data
</div>
""",
            unsafe_allow_html=True,
        )

    st.markdown("</div>", unsafe_allow_html=True)


# ============================================================
# SCREEN 2 — DOCUMENT WORKSPACE
# ============================================================
def render_workspace():
    # ---------- Sidebar ----------
    with st.sidebar:
        st.markdown(
            """
<div class="dm-workspace-brand">
  <div class="dm-small-logo">▤</div>
  <div>
    <div class="dm-workspace-name">Document Library</div>
    <div class="dm-workspace-sub">Private AI workspace</div>
  </div>
</div>
""",
            unsafe_allow_html=True,
        )

        st.markdown(
            '<p style="font-size:11px;color:#64748b;line-height:1.55;margin-top:14px">'
            "Upload files to create your private search index."
            "</p>",
            unsafe_allow_html=True,
        )

        if st.button("↔  Disconnect", use_container_width=True):
            st.session_state.connected = False
            st.session_state.client = None
            st.session_state.api_key = ""
            reset_workspace()
            st.rerun()

        st.markdown("**Upload files**")

        uploaded = st.file_uploader(
            "Add more files",
            type=["pdf", "docx", "doc", "xlsx", "xls", "csv", "pptx", "ppt", "txt", "md"],
            accept_multiple_files=True,
            label_visibility="collapsed",
        )

        if uploaded:
            existing = {d["signature"] for d in st.session_state.documents}
            new_files = []
            for f in uploaded:
                sig = sha256_bytes(f.getvalue())
                if sig not in existing:
                    new_files.append(f)

            if new_files:
                with st.spinner("Indexing documents..."):
                    try:
                        new_chunks = build_chunks(new_files)
                        if new_chunks:
                            new_vectors = get_embeddings(
                                st.session_state.client,
                                [c["text"] for c in new_chunks],
                            )
                            if st.session_state.embeddings is None:
                                st.session_state.embeddings = new_vectors
                            else:
                                st.session_state.embeddings = np.vstack(
                                    [st.session_state.embeddings, new_vectors]
                                )
                            st.session_state.chunks.extend(new_chunks)

                        for f in new_files:
                            st.session_state.documents.append(
                                {
                                    "name": f.name,
                                    "size": len(f.getvalue()),
                                    "signature": sha256_bytes(f.getvalue()),
                                    "chunks": sum(
                                        1 for c in new_chunks
                                        if c["metadata"]["filename"] == f.name
                                    ),
                                    "status": "Indexed",
                                }
                            )
                        st.rerun()
                    except Exception:
                        st.error(
                            "We couldn't process one of the files. "
                            "Make sure it isn't corrupted or password protected."
                        )

        st.markdown("---")
        st.markdown("**Settings**")

        st.session_state.settings["chunk_size"] = st.slider(
            "Chunk size", 200, 2000, st.session_state.settings["chunk_size"], 50
        )
        st.session_state.settings["chunk_overlap"] = st.slider(
            "Chunk overlap", 0, 400, st.session_state.settings["chunk_overlap"], 10
        )
        st.session_state.settings["retrieved_chunks"] = st.slider(
            "Retrieved chunks", 2, 10, st.session_state.settings["retrieved_chunks"]
        )
        st.session_state.settings["temperature"] = st.slider(
            "Temperature", 0.0, 1.0, st.session_state.settings["temperature"], 0.05
        )

        st.caption(f"Model: {st.session_state.model}")
        st.caption("Embeddings: text-embedding-3-small")

        if st.button("⌫  Clear documents", use_container_width=True):
            reset_workspace()
            st.rerun()

    # ---------- Header ----------
    st.markdown(
        """
<div class="dm-workspace-head">
  <div class="dm-workspace-brand">
    <div class="dm-small-logo">▤</div>
    <div>
      <div class="dm-workspace-name">DocuMind <span style="color:#635BFF">AI</span></div>
      <div class="dm-workspace-sub">Your documents. Now searchable with AI.</div>
    </div>
  </div>
  <div class="dm-workspace-trust">♙ Secure &nbsp; • &nbsp; Private &nbsp; • &nbsp; Powered by OpenAI</div>
</div>
""",
        unsafe_allow_html=True,
    )

    # ---------- Center ----------
    st.markdown(
        '<div style="text-align:center"><span class="dm-pill">Private Document Workspace</span></div>',
        unsafe_allow_html=True,
    )
    st.markdown(
        '<div class="dm-center-title">Talk to your <span style="color:#635BFF">documents.</span></div>',
        unsafe_allow_html=True,
    )
    st.markdown(
        '<div class="dm-center-copy">Upload files and get source-backed answers without leaving your workspace.</div>',
        unsafe_allow_html=True,
    )

    count = len(st.session_state.documents)
    chunks = len(st.session_state.chunks)
    status = "Indexed" if count else "Ready"

    st.markdown(
        f"""
<div class="dm-banner">● &nbsp; {count} file(s) indexed and ready.</div>
<div class="dm-stat-grid">
  <div class="dm-stat">
    <div class="dm-stat-icon" style="background:#eef2ff;color:#4f46e5">▤</div>
    <div><div class="dm-stat-number">{count}</div><div class="dm-stat-label">Documents</div></div>
  </div>
  <div class="dm-stat">
    <div class="dm-stat-icon" style="background:#eff6ff;color:#0ea5e9">◈</div>
    <div><div class="dm-stat-number">{chunks}</div><div class="dm-stat-label">Chunks</div></div>
  </div>
  <div class="dm-stat">
    <div class="dm-stat-icon" style="background:#ecfdf5;color:#10b981">✓</div>
    <div><div class="dm-stat-status">{status}</div><div class="dm-stat-label">Index status</div></div>
  </div>
  <div class="dm-stat">
    <div class="dm-stat-icon" style="background:#f5f3ff;color:#8b5cf6">✦</div>
    <div><div class="dm-stat-status">{st.session_state.model}</div><div class="dm-stat-label">AI model</div></div>
  </div>
</div>
""",
        unsafe_allow_html=True,
    )

    # ---------- Document list ----------
    if st.session_state.documents:
        for doc in st.session_state.documents:
            c1, c2 = st.columns([5, 1])
            with c1:
                icon = file_icon(doc["name"])
                st.markdown(
                    f"""
<div class="dm-file-row">
  <div class="dm-file-main">
    <span style="font-size:15px">{icon}</span>
    <div style="min-width:0">
      <div class="dm-file-name">{doc['name']}</div>
      <div class="dm-file-meta">{human_size(doc['size'])} • {doc['chunks']} chunks</div>
    </div>
  </div>
  <div class="dm-file-status">✓ Indexed</div>
</div>
""",
                    unsafe_allow_html=True,
                )
            with c2:
                if st.button("View", key=f"view_{doc['signature']}"):
                    st.session_state[f"show_{doc['signature']}"] = not st.session_state.get(
                        f"show_{doc['signature']}", False
                    )
            if st.session_state.get(f"show_{doc['signature']}", False):
                relevant = [
                    c for c in st.session_state.chunks
                    if c["metadata"]["signature"] == doc["signature"]
                ]
                with st.expander(f"Preview — {doc['name']}", expanded=True):
                    st.write(
                        "\n\n".join(c["text"][:500] for c in relevant[:3])
                        or "No extractable text found."
                    )

    else:
        st.markdown(
            """
<div class="dm-empty">
  <div class="dm-empty-icon">📄</div>
  <div style="font-size:13px;font-weight:700;color:#334155">No documents yet</div>
  <div style="font-size:10px;margin-top:5px">Upload your first document from the left sidebar to start asking questions.</div>
</div>
""",
            unsafe_allow_html=True,
        )

    # ---------- Chat history ----------
    if st.session_state.messages:
        for msg in st.session_state.messages:
            if msg["role"] == "user":
                st.markdown(
                    f"""
<div class="dm-question">
  <div class="dm-avatar-user">●</div>
  <div style="font-size:10px;font-weight:700;color:#334155">{msg['content']}</div>
</div>
""",
                    unsafe_allow_html=True,
                )
            else:
                st.markdown(
                    f"""
<div class="dm-answer">
  <div class="dm-answer-head"><div class="dm-avatar-ai">✦</div> DocuMind AI</div>
  <div class="dm-answer-body">{msg['content']}</div>
</div>
""",
                    unsafe_allow_html=True,
                )

    # ---------- Prompt ----------
    st.markdown('<div style="height:12px"></div>', unsafe_allow_html=True)
    question = st.chat_input("Ask a question about your documents...")

    if question:
        st.session_state.messages.append({"role": "user", "content": question})
        with st.spinner("Searching your documents..."):
            try:
                answer, sources = answer_question(question)
            except Exception:
                answer = (
                    "I couldn't complete the search right now. "
                    "Please check your API connection and try again."
                )
                sources = []

        st.session_state.messages.append({"role": "assistant", "content": answer})

        if sources:
            st.session_state.last_sources = sources
        st.rerun()

    if st.session_state.get("last_sources"):
        with st.expander("Sources", expanded=False):
            for source in st.session_state.last_sources:
                m = source["metadata"]
                location = ""
                if m.get("page"):
                    location = f"Page {m['page']}"
                elif m.get("sheet"):
                    location = f"Sheet: {m['sheet']}"
                elif m.get("slide"):
                    location = f"Slide {m['slide']}"
                st.markdown(
                    f"**{file_icon(m['filename'])} {m['filename']}** — "
                    f"{location or m['file_type']}  \n"
                    f"{m['preview'][:500]}"
                )

    st.markdown(
        '<div class="dm-disclaimer">AI responses may include mistakes. Please verify important information.</div>',
        unsafe_allow_html=True,
    )


# ============================================================
# RUN
# ============================================================
if st.session_state.connected and st.session_state.client:
    render_workspace()
else:
    render_connection()
